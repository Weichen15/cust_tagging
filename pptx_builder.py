import copy
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "template.pptx")
FONT = "微軟正黑體"

# Template slide indices (0-based)
_COVER = 0      # 封面
_NARR_CT = 3    # 標籤敘事 content
_CAT_DIV = 5    # 內容分隔頁範本 (Slide 6)
_CAT_TBL = 6    # 內容表格頁範本 (Slide 7)
_THANKS = 9     # Thank You


# ── 文字替換 ─────────────────────────────────────────────────

def _replace_in_para(para, old: str, new: str):
    """跨多個 run 做字串替換，保留第一個 run 的格式。"""
    full = "".join(r.text for r in para.runs)
    if old not in full:
        return
    new_text = full.replace(old, new)
    if para.runs:
        para.runs[0].text = new_text
        for r in para.runs[1:]:
            r.text = ""


def _replace_in_shape(shape, old: str, new: str):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        _replace_in_para(para, old, new)


# ── 表格填入 ─────────────────────────────────────────────────

def _fill_cell(cell, text: str, size: int = 11, bold: bool = False):
    """清除並填入儲存格文字，套用字型設定。"""
    cell.text = str(text)
    if not text:
        return
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold


def _fill_table(slide, stores: list):
    """將店家資料填入投影片上的表格。"""
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        tbl = shape.table
        n_data = len(tbl.rows) - 1  # 扣掉標題列
        for i in range(n_data):
            row_idx = i + 1
            if i < len(stores):
                s = stores[i]
                _fill_cell(tbl.cell(row_idx, 0), s.get("name", ""), bold=True)
                _fill_cell(tbl.cell(row_idx, 1), s.get("note", ""))
                _fill_cell(tbl.cell(row_idx, 2), s.get("source", ""))
            else:
                for col in range(3):
                    _fill_cell(tbl.cell(row_idx, col), "")
        return  # 一頁只有一個表格


# ── 投影片複製 / 刪除 ─────────────────────────────────────────

_LAYOUT_RT = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)


def _add_slide_from_xml(prs, saved_xml, src_part=None):
    """用已保存的 XML element 新增一張投影片，並複製圖片/媒體 relationships。"""
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 複製圖片 / 媒體 relationships（跳過 slideLayout）
    # 必須依照 rId 數字順序，確保 auto-assigned rId 與 shapes 引用的一致
    if src_part is not None:
        sorted_rels = sorted(
            src_part.rels.items(),
            key=lambda x: int(x[0][3:]),  # 'rId2' → 2
        )
        for rId, rel in sorted_rels:
            if rel.reltype == _LAYOUT_RT:
                continue
            try:
                if rel.is_external:
                    new_slide.part.rels.get_or_add_ext_rel(
                        rel.reltype, rel.target_ref
                    )
                else:
                    new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            except Exception:
                pass

    # 複製背景
    src_cSld = saved_xml.find(qn("p:cSld"))
    dst_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None:
        src_bg = src_cSld.find(qn("p:bg"))
        if src_bg is not None:
            dst_bg = dst_cSld.find(qn("p:bg"))
            if dst_bg is not None:
                dst_cSld.remove(dst_bg)
            dst_cSld.insert(0, copy.deepcopy(src_bg))

        # 複製所有圖形
        src_sp = src_cSld.find(qn("p:spTree"))
        dst_sp = new_slide.shapes._spTree
        for child in list(dst_sp)[2:]:
            dst_sp.remove(child)
        if src_sp is not None:
            for child in list(src_sp)[2:]:
                dst_sp.append(copy.deepcopy(child))

    return new_slide


def _delete_slide(prs, idx: int):
    """刪除指定索引的投影片。"""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[idx]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


# ── 主要函式 ─────────────────────────────────────────────────

def build_pptx(data: dict) -> BytesIO:
    prs = Presentation(TEMPLATE_PATH)

    label = data.get("label", "")
    summary = data.get("summary", "")
    categories = data.get("categories", [])

    # Slide 1 封面：替換標籤名稱
    cover = prs.slides[_COVER]
    for shape in cover.shapes:
        _replace_in_shape(shape, '"擷取標籤名稱"', label)

    # Slide 4 敘事內容：替換說明文字
    narr = prs.slides[_NARR_CT]
    for shape in narr.shapes:
        _replace_in_shape(
            shape,
            '"針對這個標籤的規畫做一個說明"',
            summary[:60],
        )

    # 在刪除前，先把範本 XML 與 part 都保存起來
    cat_div_xml  = copy.deepcopy(prs.slides[_CAT_DIV]._element)
    cat_div_part = prs.slides[_CAT_DIV].part
    cat_tbl_xml  = copy.deepcopy(prs.slides[_CAT_TBL]._element)
    cat_tbl_part = prs.slides[_CAT_TBL].part
    thanks_xml   = copy.deepcopy(prs.slides[_THANKS]._element)
    thanks_part  = prs.slides[_THANKS].part

    # 刪除 slides 5–10（從後往前，避免索引位移）
    for idx in range(10, 4, -1):
        _delete_slide(prs, idx)

    # 每個分類新增「分隔頁 + 表格頁」
    for cat in categories:
        cat_name = cat.get("category_name", "")
        reason   = cat.get("reason", "")[:80]
        stores   = cat.get("stores", [])

        # 分隔頁
        div_slide = _add_slide_from_xml(prs, cat_div_xml, cat_div_part)
        for shape in div_slide.shapes:
            _replace_in_shape(shape, '"內容1"', cat_name)
            _replace_in_shape(shape, "內容1",   cat_name)

        # 表格頁
        tbl_slide = _add_slide_from_xml(prs, cat_tbl_xml, cat_tbl_part)
        for shape in tbl_slide.shapes:
            _replace_in_shape(shape, '新增標籤 – "內容1"', f"新增標籤 – {cat_name}")
            _replace_in_shape(shape, '"內容1"', cat_name)
            _replace_in_shape(shape, "內容1",   cat_name)
            _replace_in_shape(
                shape,
                '"可以針對這個內容，做一個概述，例如找尋的原因、總結…等，讓閱聽者能一目了然"',
                reason,
            )
        _fill_table(tbl_slide, stores)

    # Thank You 頁
    _add_slide_from_xml(prs, thanks_xml, thanks_part)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
